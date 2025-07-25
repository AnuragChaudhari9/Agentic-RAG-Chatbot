import os
import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from docx import Document

class IngestionAgent:
    def __init__(self):
        pass

    def load_file(self, file_path):
        ext = file_path.split('.')[-1].lower()
        if ext == "pdf":
            return self._load_pdf(file_path)
        elif ext == "docx":
            return self._load_docx(file_path)
        elif ext == "pptx":
            return self._load_pptx(file_path)
        elif ext == "csv":
            return self._load_csv(file_path)
        elif ext in ["txt", "md"]:
            return self._load_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def _load_pdf(self, file_path):
        doc = fitz.open(file_path)
        text = "\n".join(page.get_text() for page in doc)
        return text

    def _load_docx(self, file_path):
        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

    def _load_pptx(self, file_path):
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    def _load_csv(self, file_path):
        df = pd.read_csv(file_path)
        return df.to_string()

    def _load_text(self, file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
